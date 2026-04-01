import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import os
from dotenv import load_dotenv

# --- 1. INITIALIZE SESSION STATES ---
load_dotenv()
api_key = os.getenv("GEMINI_API_KEY")

if 'mastery' not in st.session_state:
    st.session_state.mastery = {"Reading": 0, "Practice": 0, "Quiz": 0}
if 'quiz_data' not in st.session_state:
    st.session_state.quiz_data = []
if 'score' not in st.session_state:
    st.session_state.score = 0
if 'current_q_idx' not in st.session_state:
    st.session_state.current_q_idx = 0

# --- 2. API CONFIG ---
if api_key:
    genai.configure(api_key=api_key)
else:
    st.error("API Key not found! Please check your .env file.")

model = genai.GenerativeModel('gemini-2.5-flash')

# --- 3. PAGE CONFIG ---
st.set_page_config(page_title="AI RAG Learning Assistant", layout="wide")

# --- 4. HELPER FUNCTIONS ---
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def get_ai_response(context, user_query, level):
    source_text = context[:5000] if context else "No study material uploaded yet."
    prompt = f"Expert CS Tutor. Context: {source_text}. Question: {user_query}. Level: {level}."
    response = model.generate_content(prompt)
    return response.text

def generate_full_quiz(context):
    if not context:
        return []
    prompt = f"""
    Based on this text: {context[:4000]}
    Generate 5 Multiple Choice Questions. 
    Return them strictly in this format:
    Q: [Question] | A: [Opt1] | B: [Opt2] | C: [Opt3] | D: [Opt4] | Correct: [Letter] | Hint: [One sentence hint]
    """
    response = model.generate_content(prompt).text
    questions = []
    for line in response.strip().split('\n'):
        if "|" in line:
            parts = line.split("|")
            try:
                questions.append({
                    "q": parts[0].replace("Q:", "").strip(),
                    "options": [parts[1].strip(), parts[2].strip(), parts[3].strip(), parts[4].strip()],
                    "correct": parts[5].replace("Correct:", "").strip(),
                    "hint": parts[6].replace("Hint:", "").strip()
                })
            except:
                continue
    return questions

# --- 5. SIDEBAR DASHBOARD ---
st.sidebar.title("🛤️ Learning Dashboard")

uploaded_file = st.sidebar.file_uploader("📂 Upload Study Material", type=["pdf", "pptx"])

# Define context_data properly
context_data = ""
if uploaded_file:
    if uploaded_file.name.endswith(".pdf"):
        context_data = extract_text_from_pdf(uploaded_file)
    else:
        context_data = extract_text_from_pptx(uploaded_file)
    st.sidebar.success(f"Loaded: {uploaded_file.name}")

# Progress Metrics
st.sidebar.metric("📖 Reading", f"{st.session_state.mastery['Reading']}%")
st.sidebar.metric("📝 Practice", f"{st.session_state.mastery['Practice']}%")
st.sidebar.metric("🧠 Quiz", f"{st.session_state.mastery['Quiz']}%")

if st.sidebar.button("🔄 Reset Progress"):
    st.session_state.mastery = {"Reading": 0, "Practice": 0, "Quiz": 0}
    st.rerun()

# --- 6. MAIN UI WITH TABS ---
st.markdown("<h1 style='text-align: center; color: #4CAF50;'>🤖 AI Smart Learning Assistant</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center; color: gray;'>Your personalized learning path powered by RAG</p>", unsafe_allow_html=True)

tab1, tab2, tab3 = st.tabs(["📖 Learn", "📝 Practice", "🧠 Quiz"])

with tab1:
    st.header("📖 Learn")
    st.subheader("Explain a Topic")
    topic_input = st.text_input("What do you want to learn today?")
    level = st.select_slider("Level", options=["Summary", "Standard", "Deep Dive"])
    if st.button("🔍 Explain Topic"):
        if topic_input:
            st.session_state.last_explanation = get_ai_response(context_data, topic_input, level)
            st.session_state.mastery['Reading'] = min(100, st.session_state.mastery['Reading'] + 20)
        else:
            st.error("Please enter a topic!")
    if 'last_explanation' in st.session_state:
        with st.expander("📖 Explanation", expanded=True):
            st.write(st.session_state.last_explanation)

with tab2:
    st.header("📝 Practice")
    st.subheader("Generate Practice Summary")
    if st.button("⚡ Generate Summary"):
        if context_data:
            practice_prompt = f"Summarize 5 key terms from: {context_data[:3000]}"
            st.session_state.last_practice = model.generate_content(practice_prompt).text
            st.session_state.mastery['Practice'] = min(100, st.session_state.mastery['Practice'] + 20)
        else:
            st.error("Please upload a file first!")
    if 'last_practice' in st.session_state:
        with st.expander("📝 Practice Notes", expanded=True):
            st.success(st.session_state.last_practice)

with tab3:
    st.header("🧠 Quiz")
    st.subheader("Final Assessment")
    if st.button("🏁 Generate Quiz"):
        if context_data:
            st.session_state.quiz_data = generate_full_quiz(context_data)
            st.session_state.current_q_idx = 0
            st.session_state.score = 0
        else:
            st.error("Please upload a file first!")

    if st.session_state.quiz_data:
        idx = st.session_state.current_q_idx
        if idx < 5:
            q = st.session_state.quiz_data[idx]
            st.write(f"**Q{idx+1}:** {q['q']}")

            # Use a form so buttons don't clash
            with st.form(key=f"quiz_form_{idx}"):
                user_answer = st.radio("Choose:", q["options"])
                submitted = st.form_submit_button("Submit Answer")
                hint_btn = st.form_submit_button("💡 Get a Hint")

                if submitted:
                    if user_answer.startswith(q["correct"]):
                        st.success("✅ Correct!")
                        st.session_state.score += 1
                        # Increment mastery gradually
                        st.session_state.mastery['Quiz'] = min(100, st.session_state.mastery['Quiz'] + 5)
                    else:
                        st.error(f"❌ Wrong! Correct: {q['correct']}")

                    st.session_state.current_q_idx += 1
                    st.rerun()

                if hint_btn:
                    st.info(f"Hint: {q['hint']}")

        else:
            st.write(f"🎉 Quiz Finished! Score: {st.session_state.score}/5")
            if st.button("Start New Quiz"):
                st.session_state.quiz_data = []
                st.rerun()
